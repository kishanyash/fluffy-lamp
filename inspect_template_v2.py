
from pptx import Presentation
import os

def inspect():
    with open("placeholders.txt", "w", encoding="utf-8") as f:
        path = "master_template.pptx"
        if not os.path.exists(path):
            f.write("ERROR: master_template.pptx not found")
            return

        try:
            prs = Presentation(path)
            f.write(f"Template Analysis for: {path}\n")
            f.write(f"Total Slides: {len(prs.slides)}\n\n")

            all_placeholders = set()

            for i, slide in enumerate(prs.slides):
                f.write(f"--- Slide {i+1} ---\n")
                found_on_slide = []
                
                # Check shapes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        if "{{" in text:
                            found_on_slide.append(f"TEXT: {text.strip()}")
                            # Extract placeholder names
                            import re
                            matches = re.findall(r"\{\{(.*?)\}\}", text)
                            all_placeholders.update(matches)
                            
                    if shape.has_table:
                        # Check first cell for table placeholder
                        try:
                            first_cell = shape.table.cell(0,0).text_frame.text.strip()
                            if "{{" in first_cell:
                                found_on_slide.append(f"TABLE: {first_cell}")
                                matches = re.findall(r"\{\{(.*?)\}\}", first_cell)
                                all_placeholders.update(matches)
                        except:
                            pass

                if not found_on_slide:
                    f.write("(No placeholders found)\n")
                else:
                    for item in found_on_slide:
                        f.write(f"  {item}\n")
                f.write("\n")

            f.write("\n=== SUMMARY OF PLACEHOLDERS FOUND ===\n")
            for p in sorted(all_placeholders):
                f.write(f"- {p}\n")

        except Exception as e:
            f.write(f"ERROR: {e}")

if __name__ == "__main__":
    inspect()
