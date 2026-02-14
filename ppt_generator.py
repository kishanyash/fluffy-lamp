"""
PPT Generator Script for Research Reports
==========================================
This script automates the generation of PowerPoint presentations
from Supabase data received via n8n webhook.

Template Structure (11 slides):
- Slide 1: Title (company_name, nse_symbol, bom_code, rating)
- Slide 2: Company Background
- Slide 3: Business Model
- Slide 4: Management Analysis
- Slide 5: Industry Overview
- Slide 6: Key Industry Tailwinds
- Slide 7: Demand Drivers
- Slide 8: Industry Risks
- Slide 9: Financials (summary_table text + 4 chart quadrants)
- Slide 10: Summary Charts (summary_charts text)
- Slide 11: Disclaimers
"""

import os
import re
import requests
from io import BytesIO
from datetime import datetime
from typing import Dict, Optional, Any, List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PPTGenerator:
    """
    A class to generate PowerPoint presentations from research report data.
    """

    # Chart image positions (slide index is 0-based)
    # Note: Slide 9 (index 8) has 4 financial chart quadrants AND the summary_table image
    # Slide 10 (index 9) has the chart_custom image
    CHART_POSITIONS = {
        'chart_profit_loss': {
            'slide': 8,  # Slide 9
            'position': {'left': 0.16, 'top': 1.4, 'width': 4.8, 'height': 2.2}
        },
        'chart_balance_sheet': {
            'slide': 8,  # Slide 9
            'position': {'left': 5.15, 'top': 1.4, 'width': 4.8, 'height': 2.2}
        },
        'chart_cash_flow': {
            'slide': 8,  # Slide 9
            'position': {'left': 0.10, 'top': 4.3, 'width': 4.8, 'height': 2.2}
        },
        'chart_ratio_analysis': {
            'slide': 8,  # Slide 9
            'position': {'left': 5.10, 'top': 4.3, 'width': 4.8, 'height': 2.2}
        },
        'summary_table': {
            'slide': 8,  # Slide 9 (Summary in Tables)
            'position': {'left': 0.5, 'top': 0.75, 'width': 9.0, 'height': 4.5}
        },
        'chart_custom': {
            'slide': 9,  # Slide 10 (Summary in Charts)
            'position': {'left': 0.5, 'top': 0.75, 'width': 9.0, 'height': 4.5}
        },
        'price_chart': {
             'slide': 0, # Slide 1 (Title)
             'position': {'left': 7.0, 'top': 2.0, 'width': 3.0, 'height': 2.0}
        },
    }

    def __init__(self, template_path: str):
        """Initialize the PPT Generator with a template."""
        self.template_path = template_path
        self.prs = None

    def load_template(self) -> None:
        """Load the PowerPoint template."""
        if not os.path.exists(self.template_path):
            raise FileNotFoundError(f"Template not found: {self.template_path}")
        self.prs = Presentation(self.template_path)
        print(f"  Loaded template with {len(self.prs.slides)} slides")

    def parse_markdown_to_text(self, markdown_text: str) -> str:
        """
        Convert markdown text to clean plain text.
        Preserves paragraph structure but removes markdown formatting.
        """
        if not markdown_text:
            return ""

        text = markdown_text

        # Remove markdown headers but keep the text
        text = re.sub(r'^#{1,6}\s*(.+)$', r'\1', text, flags=re.MULTILINE)

        # Convert bold/italic markers
        text = re.sub(r'\*\*\*(.+?)\*\*\*', r'\1', text)  # Bold italic
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)      # Bold
        text = re.sub(r'(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)', r'\1', text)  # Italic
        text = re.sub(r'__(.+?)__', r'\1', text)          # Bold alt
        text = re.sub(r'(?<!_)_(?!_)(.+?)(?<!_)_(?!_)', r'\1', text)  # Italic alt

        # Remove link formatting but keep text
        text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)

        # Clean up excessive newlines (keep double newlines for paragraphs)
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # Remove leading/trailing whitespace from each line
        lines = [line.strip() for line in text.split('\n')]
        text = '\n'.join(lines)

        return text.strip()

    def download_image(self, url: str) -> Optional[BytesIO]:
        """Download an image from URL and return as BytesIO object."""
        if not url or url in ("[null]", "null", None, ""):
            return None

        try:
            print(f"    Downloading: {url[:60]}...")
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            image_data = BytesIO(response.content)
            image_data.seek(0)
            return image_data
        except Exception as e:
            print(f"    Error downloading image: {e}")
            return None

    def find_shape_with_placeholder(self, placeholder_name: str):
        """
        Find the shape containing the placeholder text.
        Returns (slide, shape) tuple or (None, None) if not found.
        """
        placeholder_pattern = f"{{{{{placeholder_name}}}}}"
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                # Check full text of the shape
                full_text = ""
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        full_text += run.text
                
                if placeholder_pattern in full_text:
                    return slide, shape
        
        return None, None

    def replace_shape_text(self, shape, new_text: str, font_size: int = 10) -> bool:
        """
        Replace the entire text content of a shape with new text.
        Properly handles text frame formatting to prevent overflow and overlapping.
        """
        if not shape.has_text_frame:
            return False
        
        tf = shape.text_frame
        
        # Set text frame properties to prevent overflow
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Auto-shrink text to fit
        
        # Set margins (Inches)
        tf.margin_top = Inches(0.25)  # Reduced from 0.5 to move text up
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_bottom = Inches(0.1)
        
        # Set vertical anchor to top
        tf.vertical_anchor = MSO_ANCHOR.TOP
        
        # Clear existing content and set new text
        # First paragraph
        if not tf.paragraphs:
            tf.add_paragraph()
            
        p = tf.paragraphs[0]
        p.clear()
        
        # Clean text
        clean_text = new_text.strip()
        
        # Add text run
        run = p.add_run()
        run.text = clean_text
        
        # Set font properties - let auto-fit handle size usually, but set a max starting point
        run.font.size = Pt(float(font_size)) 
        run.font.name = "Arial"
        
        # Remove any extra paragraphs
        while len(tf.paragraphs) > 1:
            # Can't directly remove paragraphs easily in some python-pptx versions
            # So we clear them
            for para in tf.paragraphs[1:]:
                para.clear()
            break
        
        return True

    def find_and_replace_placeholder(self, placeholder_name: str, new_text: str, font_size: int = 10) -> int:
        """
        Find and replace {{placeholder_name}} with new text.
        Uses proper text replacement to avoid overflow issues.
        """
        placeholder_pattern = f"{{{{{placeholder_name}}}}}"
        replacements = 0

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                tf = shape.text_frame
                
                # Get full text of the shape
                full_text = ""
                for para in tf.paragraphs:
                    full_text += ''.join(run.text for run in para.runs)
                
                # Check if placeholder exists
                if placeholder_pattern not in full_text:
                    continue
                
                # For simple single-placeholder shapes, replace entire content
                if full_text.strip() == placeholder_pattern:
                    # This is a simple placeholder-only shape
                    self.replace_shape_text(shape, new_text, font_size)
                    replacements += 1
                else:
                    # Multiple placeholders or mixed content - do inline replacement
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if placeholder_pattern in run.text:
                                run.text = run.text.replace(placeholder_pattern, new_text)
                                if font_size:
                                    run.font.size = Pt(font_size)
                                replacements += 1
                        
                        # Also try combined runs
                        combined = ''.join(run.text for run in para.runs)
                        if placeholder_pattern in combined and replacements == 0:
                            new_combined = combined.replace(placeholder_pattern, new_text)
                            if para.runs:
                                para.runs[0].text = new_combined
                                if font_size:
                                    para.runs[0].font.size = Pt(font_size)
                                for run in para.runs[1:]:
                                    run.text = ""
                                replacements += 1

        return replacements

    def parse_markdown_table_to_data(self, markdown_text: str) -> List[List[str]]:
        """
        Parses a markdown table into a list of lists (rows of columns).
        Example input:
        | Header 1 | Header 2 |
        |---|---|
        | Row 1 Col 1 | Row 1 Col 2 |
        """
        if not markdown_text:
            return []
            
        lines = markdown_text.strip().split('\n')
        table_data = []
        
        for line in lines:
            # Skip separator lines (e.g. |---|---|)
            if '---' in line:
                continue
            # Skip empty lines
            if not line.strip():
                continue
            # Check if likely a row
            if '|' in line:
                # Split by pipe, strip whitespace
                row = [cell.strip() for cell in line.split('|')]
                # Filter out empty strings from leading/trailing pipes
                row = [cell for cell in row if cell]
                if row:
                    table_data.append(row)
                    
        return table_data

    def populate_table_shape(self, shape, data: List[List[str]], font_size: int = 10):
        """
        Populate a PowerPoint table shape with data (list of lists).
        """
        if not shape.has_table:
            return
            
        table = shape.table
        
        # Iterate over data rows
        for r_idx, row_data in enumerate(data):
            # If we run out of table rows, stop (or add rows if supported/needed)
            if r_idx >= len(table.rows):
                break
                
            for c_idx, cell_value in enumerate(row_data):
                # If we run out of table cols, stop
                if c_idx >= len(table.columns):
                    break
                    
                cell = table.cell(r_idx, c_idx)
                # Use same text replacement logic to handle formating
                self.replace_shape_text(cell, str(cell_value), font_size)

    def find_and_populate_table(self, placeholder_text: str, data: List[List[str]], font_size: int = 10) -> bool:
        """
        Find a table that contains the specific placeholder in its first cell (0,0)
        and populate it with the provided data.
        """
        target_pattern = f"{{{{{placeholder_text}}}}}"
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    # Check first cell for placeholder
                    try:
                        first_cell_text = shape.table.cell(0, 0).text_frame.text.strip()
                        if target_pattern in first_cell_text:
                            # Clear the placeholder logic from the first cell effectively by overwriting 
                            # when populating, or we treat the headers as row 0.
                            
                            print(f"    -> Found table with placeholder '{placeholder_text}' on Slide {self.prs.slides.index(slide)+1}")
                            self.populate_table_shape(shape, data, font_size)
                            return True
                    except Exception:
                        continue
        return False

    def add_image_to_slide(self, slide_idx: int, image_data: BytesIO,
                           left: float, top: float,
                           width: float, height: Optional[float] = None) -> bool:
        """Add an image to a specific slide."""
        if slide_idx >= len(self.prs.slides):
            print(f"    Warning: Slide {slide_idx + 1} does not exist")
            return False

        try:
            slide = self.prs.slides[slide_idx]
            image_data.seek(0)

            if height:
                slide.shapes.add_picture(
                    image_data, 
                    Inches(left), Inches(top),
                    width=Inches(width), height=Inches(height)
                )
            else:
                slide.shapes.add_picture(
                    image_data, 
                    Inches(left), Inches(top),
                    width=Inches(width)
                )

            return True
        except Exception as e:
            print(f"    Error adding image: {e}")
            return False

    def add_debug_grid(self, slide_idx: int):
        """Add visual debug lines to the slide."""
        try:
            slide = self.prs.slides[slide_idx]
            from pptx.util import Inches
            from pptx.enum.shapes import MSO_CONNECTOR
            from pptx.dml.color import RGBColor

            # Draw Red Line at Top = 1.1 inches (Target top)
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(1.1), Inches(10), Inches(1.1)
            )
            line.line.color.rgb = RGBColor(255, 0, 0)
            line.line.width = Inches(0.05)

            # Draw Green Line at Top = 6.6 inches (Target bottom approx)
            line2 = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(6.6), Inches(10), Inches(6.6)
            )
            line2.line.color.rgb = RGBColor(0, 255, 0)
            line2.line.width = Inches(0.05)
            
            print(f"    DEBUG: Added red/green lines to Slide {slide_idx+1}")
            
            # Print slide dimensions
            print(f"    DEBUG: Slide width={self.prs.slide_width/914400} inches, height={self.prs.slide_height/914400} inches")

        except Exception as e:
            # Typically imports might fail if python-pptx version is old or structured differently
            # We try to import inside to be safe or just print error
            print(f"    Debug error: {e}")

    def calculate_font_size(self, text: str, max_chars: int = 2000) -> float:
        """
        Calculate appropriate font size based on text length.
        Longer text gets smaller font.
        """
        text_len = len(text)
        
        if text_len < 500:
            return 11.0  # Standard body text
        elif text_len < 1000:
            return 10.5
        elif text_len < 1500:
            return 10.0
        elif text_len < 2000:
            return 9.0
        elif text_len < 3000:
            return 8.0 # Let auto-fit shrink further if needed
        else:
            return 7.0 # Set a readable base even for long text

    def fetch_bom_code(self, symbol: str, company_name: str) -> str:
        """
        Attempt to fetch BOM code from Yahoo Finance search API.
        Searches for the symbol or company name and looks for a result ending in '.BO'.
        Falls back to a known list of common Indian stocks.
        """
        # Fallback list of common Indian stocks (NSE symbol -> BSE code)
        KNOWN_BSE_CODES = {
            'WIPRO': '507685',
            'TCS': '532540',
            'INFY': '500209',
            'RELIANCE': '500325',
            'HDFCBANK': '500180',
            'ICICIBANK': '532174',
            'SBIN': '500112',
            'BHARTIARTL': '532454',
            'ITC': '500875',
            'HINDUNILVR': '500696',
            'KOTAKBANK': '500247',
            'LT': '500510',
            'AXISBANK': '532215',
            'ASIANPAINT': '500820',
            'MARUTI': '532500',
            'TATAMOTORS': '500570',
            'SUNPHARMA': '524715',
            'TITAN': '500114',
            'BAJFINANCE': '500034',
            'HCLTECH': '532281',
        }
        
        # Check fallback list first
        symbol_upper = symbol.upper() if symbol else ''
        if symbol_upper in KNOWN_BSE_CODES:
            print(f"    -> Found BSE code in fallback list: {KNOWN_BSE_CODES[symbol_upper]}")
            return KNOWN_BSE_CODES[symbol_upper]
        
        try:
            # Try searching by symbol first, then company name
            queries = [q for q in [symbol, company_name] if q]
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }

            for query in queries:
                print(f"    -> Searching Yahoo Finance for: {query}")
                url = f"https://query2.finance.yahoo.com/v1/finance/search?q={query}&quotesCount=10&newsCount=0"
                response = requests.get(url, headers=headers, timeout=10)
                
                print(f"    -> Response status: {response.status_code}")
                
                if response.status_code == 200:
                    data = response.json()
                    quotes = data.get('quotes', [])
                    print(f"    -> Found {len(quotes)} quotes")
                    
                    for quote in quotes:
                        symbol_ticker = quote.get('symbol', '')
                        # Look for BSE symbols (usually end with .BO)
                        if symbol_ticker.endswith('.BO'):
                            bse_code = symbol_ticker.split('.')[0]
                            print(f"    -> Found BSE symbol: {symbol_ticker} -> {bse_code}")
                            return bse_code
                            
        except Exception as e:
            print(f"    Warning: Could not fetch BOM code: {e}")
            
        print("    -> No BSE code found")
        return ' '

    def populate_from_data(self, data: Dict[str, Any]) -> Dict[str, bool]:
        """
        Populate the presentation with data from the Supabase record.
        """
        print("\n" + "=" * 60)
        print("POPULATING PRESENTATION")
        print("=" * 60)

        results = {}

        # ===== TABLE POPULATION =====
        print("\n--- Table Population ---")
        
        # 1. Try to use explicit markdown if provided (highest priority)
        # 1. Try to use explicit markdown if provided (highest priority)
        financial_val = data.get('financial_performance', '')
        table_data = []
        financial_text_summary = ""
        
        if financial_val and '|' in str(financial_val):
            print("  Found markdown table in 'financial_performance'. Parsing...")
            table_data = self.parse_markdown_table_to_data(str(financial_val))
        else:
            print("  'financial_performance' appears to be text summary.")
            financial_text_summary = str(financial_val) if financial_val else ""
        
        # 2. If no markdown, try to construct from individual DB fields (equity_universe data)
        elif 'revenue_fy2024' in data or 'revenue_ttm' in data:
            print("  Constructing table from equity_universe fields...")
            
            # Helper to safely get numeric value formatted
            def get_val(key, fmt="{:,.0f}"):
                val = data.get(key)
                if val is None: return "-"
                try:
                    return fmt.format(float(val))
                except:
                    return str(val)

            # Define the structure based on your screenshot
            # Header Row
            headers = ["Particulars", "FY24A", "FY25A", "FY26E", "FY27E", "FY28E"]
            
            # Data Rows
            rows = [
                # Sales
                ["Sales", 
                 get_val('revenue_fy2024'), get_val('revenue_fy2025'), 
                 get_val('revenue_fy2026e'), get_val('revenue_fy2027e'), get_val('revenue_fy2028e')],
                
                # Sales Growth (YoY %) - You might need to calculate this if not in DB
                ["YoY% growth", 
                 get_val('sales_growth_yoy_qtr', "{:.1f}"), "-", "-", "-", "-"], 

                # EBITDA
                ["EBITDA", 
                 get_val('ebitda_fy2024'), get_val('ebitda_fy2025'), 
                 get_val('ebitda_fy2026e'), get_val('ebitda_fy2027e'), get_val('ebitda_fy2028e')],

                # EBITDA Margin (%)
                ["% Margin", 
                 get_val('ebitda_margin_fy2024', "{:.1f}"), get_val('ebitda_margin_fy2025', "{:.1f}"), 
                 get_val('ebitda_margin_fy2026e', "{:.1f}"), get_val('ebitda_margin_fy2027e', "{:.1f}"), get_val('ebitda_margin_fy2028e', "{:.1f}")],

                # PAT
                ["PAT", 
                 get_val('pat_fy2024'), get_val('pat_fy2025'), 
                 get_val('pat_fy2026e'), get_val('pat_fy2027e'), get_val('pat_fy2028e')],

                # PAT Growth
                ["YoY% growth", 
                 get_val('pat_growth_qoq', "{:.1f}"), "-", "-", "-", "-"], # Using QoQ as placeholder if YoY missing

                # P/E
                ["P/E", 
                 get_val('pe_ttm', "{:.1f}"), get_val('pe_fy2025', "{:.1f}"), # Note: pe_fy2025 might not exist, check DB keys
                 get_val('pe_fy2026e', "{:.1f}"), get_val('pe_fy2027e', "{:.1f}"), get_val('pe_fy2028e', "{:.1f}")],

                # P/B (Book Value) - We have book_value, need P/B calculation or field
                ["P/B", 
                 "-", "-", "-", "-", "-"] 
            ]
            
            table_data = [headers] + rows
        
        # 3. Populate if we have data
        if table_data:
            success = self.find_and_populate_table('financial_table', table_data, font_size=10)
            print(f"  Financial Table: {'[OK] Populated' if success else '[FAILED] Table placeholder {{financial_table}} not found'}")
        else:
            print("  Financial Table: No data found (markdown or DB fields)")
        
        # ===== TEXT REPLACEMENTS =====
        print("\n--- Text Replacements ---")
        
        # Get or fetch BOM code (must be numeric like "507685")
        bom_code = data.get('bom_code', '')
        # Check if bom_code is valid (should be numeric)
        is_valid_bom = bom_code and str(bom_code).strip().isdigit()
        
        if not is_valid_bom:
            print(f"  BOM Code '{bom_code}' is invalid (not numeric). Fetching from Yahoo Finance...")
            symbol = data.get('nse_symbol', data.get('symbol', ''))
            name = data.get('company_name', '')
            bom_code = self.fetch_bom_code(symbol, name)
            print(f"  -> Found: {bom_code}" if bom_code.strip() else "  -> Not found")
        else:
            print(f"  BOM Code: {bom_code} (provided)")
        
        # Get rating, default to N/A if missing
        rating = data.get('rating', '')
        if not rating or str(rating).strip() == '':
            rating = 'N/A'
        print(f"  DEBUG: Rating/Recommendation value: '{rating}'")
        
        # Define placeholder mappings with their data sources
        text_mappings = [
            ('company_name', data.get('company_name', ''), 24),  # Large font for title
            ('nse_symbol', data.get('nse_symbol', data.get('symbol', '')), 14),
            ('bom_code', bom_code, 14),
            ('recommendation', rating, 14),
            ('today_date', data.get('today_date', datetime.now().strftime('%Y-%m-%d')), 14),
            ('company_background', self.parse_markdown_to_text(data.get('company_background', '')), None),
            ('business_model', self.parse_markdown_to_text(data.get('business_model', '')), None),
            ('management_analysis', self.parse_markdown_to_text(data.get('management_analysis', '')), None),
            ('industry_overview', self.parse_markdown_to_text(data.get('industry_overview', '')), None),
            ('industry_tailwinds', self.parse_markdown_to_text(data.get('industry_tailwinds', data.get('key_industry', ''))), None),
            ('demand_drivers', self.parse_markdown_to_text(data.get('demand_drivers', '')), None),
            ('industry_risk', self.parse_markdown_to_text(data.get('industry_risks', data.get('industry_risk', ''))), None),
            
            # --- NEW FIELDS ---
            ('market_positioning', self.parse_markdown_to_text(data.get('market_positioning', '')), None),
            ('financial_performance', financial_text_summary, None),
            ('growth_outlook', self.parse_markdown_to_text(data.get('growth_outlook', '')), None),
            ('valuation_recommendation', self.parse_markdown_to_text(data.get('valuation_recommendation', '')), None),
            ('key_risks', self.parse_markdown_to_text(data.get('key_risks', '')), None),
            ('company_insider', self.parse_markdown_to_text(data.get('company_insider', '')), None),
            
            # Scripts
            ('podcast_script', self.parse_markdown_to_text(data.get('podcast_script', '')), None),
            ('video_script', self.parse_markdown_to_text(data.get('video_script', '')), None),

            # Clear these placeholders as they will be replaced by images
            ('summary_table', ' ', None),  
            ('chart_custom', ' ', None),
            ('price_chart', ' ', None),
        ]

        for placeholder, value, fixed_font_size in text_mappings:
            if value:
                # Calculate font size if not fixed
                if fixed_font_size:
                    font_size = fixed_font_size
                else:
                    font_size = self.calculate_font_size(value)
                
                # Limit text length to prevent extreme overflow - REMOVED LIMIT
                # max_chars = 3000
                # if len(value) > max_chars:
                #    value = value[:max_chars] + "\n\n[Content truncated...]"
                
                count = self.find_and_replace_placeholder(placeholder, value, font_size)
                results[placeholder] = count > 0
                
                char_info = f"{len(value)} chars, {font_size}pt"
                status = f"[OK] Replaced ({char_info})" if count > 0 else "[MISSING] Placeholder not found"
                print(f"  {placeholder}: {status}")
            else:
                results[placeholder] = False
                print(f"  {placeholder}: [MISSING] No data provided")

        # ===== IMAGE INSERTIONS =====
        print("\n--- Image Insertions ---")
        
        image_fields = {
            'chart_profit_loss': data.get('chart_profit_loss'),
            'chart_balance_sheet': data.get('chart_balance_sheet'),
            'chart_cash_flow': data.get('chart_cash_flow'),
            'chart_ratio_analysis': data.get('chart_ratio_analysis'),
            'summary_table': data.get('summary_table'),  # Slide 9 Image
            'chart_custom': data.get('chart_custom'),    # Slide 10 Image
            'price_chart': data.get('price_chart'),      # New chart
        }

        for field_name, url in image_fields.items():
            if url and url not in ("[null]", "null", None, ""):
                config = self.CHART_POSITIONS.get(field_name)
                if config:
                    print(f"  {field_name}:")
                    image_data = self.download_image(url)
                    if image_data:
                        pos = config.get('position', {})
                        success = self.add_image_to_slide(
                            config['slide'],
                            image_data,
                            left=pos.get('left', 1.0),
                            top=pos.get('top', 1.5),
                            width=pos.get('width', 8.0),
                            height=pos.get('height')
                        )
                        results[field_name] = success
                        print(f"    -> Slide {config['slide'] + 1}: {'[OK] Added' if success else '[FAILED]'}")
                    else:
                        results[field_name] = False
                        print(f"    -> [FAILED] Download failed")
            else:
                results[field_name] = False
                print(f"  {field_name}: [MISSING] No URL provided")

        return results

    def save(self, output_path: str) -> str:
        """Save the presentation to a file."""
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        self.prs.save(output_path)
        print(f"\n[OK] Presentation saved to: {output_path}")
        return output_path


def generate_report_ppt(data: Dict[str, Any], 
                        template_path: str,
                        output_dir: str = "./output") -> str:
    """
    Main function to generate a PowerPoint report.
    """
    # Generate output filename
    report_id = data.get('report_id', 'unknown')
    symbol = data.get('symbol', data.get('nse_symbol', 'report'))
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    # Clean report_id for filename
    report_id_clean = report_id[:8] if len(report_id) > 8 else report_id
    output_filename = f"{symbol}_{report_id_clean}_{timestamp}.pptx"
    output_path = os.path.join(output_dir, output_filename)

    # Create generator and process
    generator = PPTGenerator(template_path)
    generator.load_template()
    results = generator.populate_from_data(data)
    generator.save(output_path)

    # Summary
    successful = sum(1 for v in results.values() if v)
    total = len(results)
    print(f"\n{'=' * 60}")
    print(f"SUMMARY: {successful}/{total} fields processed successfully")
    print(f"{'=' * 60}")

    return output_path


# ============================================================
# EXAMPLE USAGE AND TESTING
# ============================================================
if __name__ == "__main__":
    # Example data structure (as received from n8n/Supabase)
    example_data = {
        "report_id": "c49b2aa1-80eb-4436-b14c-2a74d7966feb",
        "company_name": "Vedanta Ltd.",
        "nse_symbol": "VEDL",
        "bom_code": "500295",
        "rating": "BUY",
        "today_date": "2026-02-07",
        "company_background": """Company Background

Vedanta Limited is a globally diversified natural resources company with business operations in India, South Africa, Namibia, and Australia. The company is one of the world's largest diversified natural resources companies.

Key Business Segments:
• Zinc Business: One of the largest integrated producers of zinc-lead
• Aluminum Business: India's largest aluminum producer
• Oil and Gas: Significant crude oil producer in India
• Iron Ore: Mining operations in Goa and Karnataka
• Copper: Copper smelting and refining operations

History and Evolution:
Founded in 1976, Vedanta has grown through strategic acquisitions and organic expansion. The company was originally focused on mining and has diversified into various natural resources over the decades.

Market Position:
Vedanta holds leadership positions in multiple segments of the Indian natural resources industry, with significant global presence in key commodities.""",

        "business_model": """Business Model Explanation

Revenue Streams:
Vedanta generates revenue through multiple integrated business segments including mining operations, smelting and refining, and oil and gas production.

1. Mining Operations
• Extraction of zinc, lead, silver, iron ore
• Open-pit and underground mining operations
• Mineral processing and concentration

2. Smelting and Refining
• Aluminum smelting operations
• Copper cathode production
• Zinc and lead refining

3. Oil and Gas Production
• Crude oil extraction from Rajasthan fields
• Natural gas production

Value Chain Integration:
The company maintains vertical integration across exploration, mining, processing, and marketing. This integration provides cost advantages and supply chain control.

Key Competitive Advantages:
• Low-cost production capabilities
• Diverse commodity portfolio reducing risk
• Strong operational expertise
• Strategic asset locations""",

        "management_analysis": """Management Analysis

Leadership Team:
Anil Agarwal - Chairman: Founder and visionary leader with over 40 years of industry experience, known for bold strategic decisions.

Key Management Metrics:
• Experience: Excellent
• Track Record: Strong
• Corporate Governance: Good
• Capital Allocation: Above Average

Strategic Direction:
The management has outlined a clear growth strategy focusing on capacity expansion in aluminum and zinc, exploration and development of new resources, ESG improvements and sustainability initiatives, and digital transformation of operations.""",

        "industry_overview": """Industry Overview

Industry Size & Structure:
The mining and metals industry is a significant contributor to the global economy. The Total Addressable Market for this sector is vast, driven by demand for essential metals such as aluminum, copper, zinc, and iron ore.

Market Dynamics:
• Total global mining market: $2.1 trillion
• Base metals segment: $650 billion
• Expected CAGR: 4.5% (2024-2030)

Indian Market Position:
• India is the 3rd largest producer of coal
• 4th largest producer of iron ore
• Significant growth potential in base metals""",

        # --- MOCK FINANCIAL DATA (equity_universe fields) ---
        "revenue_fy2024": 150000, "revenue_fy2025": 165000, "revenue_fy2026e": 180000, "revenue_fy2027e": 200000, "revenue_fy2028e": 225000,
        "sales_growth_yoy_qtr": 12.5,
        "ebitda_fy2024": 45000, "ebitda_fy2025": 50000, "ebitda_fy2026e": 55000, "ebitda_fy2027e": 62000, "ebitda_fy2028e": 70000,
        "ebitda_margin_fy2024": 30.0, "ebitda_margin_fy2025": 30.3, "ebitda_margin_fy2026e": 30.5, "ebitda_margin_fy2027e": 31.0, "ebitda_margin_fy2028e": 31.1,
        "pat_fy2024": 12000, "pat_fy2025": 14000, "pat_fy2026e": 16000, "pat_fy2027e": 19000, "pat_fy2028e": 23000,
        "pat_growth_qoq": 15.2,
        "pe_ttm": 15.4, "pe_fy2025": 14.2, "pe_fy2026e": 12.5, "pe_fy2027e": 10.8, "pe_fy2028e": 9.2,

        "industry_tailwinds": """Key Industry Tailwinds

Structural Growth Drivers:

1. Infrastructure Development
• Government's infrastructure push (PM Gati Shakti)
• National Infrastructure Pipeline: ₹111 lakh crore investment
• Increased demand for steel, aluminum, and copper

2. Electric Vehicle Revolution
• EV adoption driving copper and aluminum demand
• Battery metals gaining importance
• India's EV sales growing at 40%+ CAGR

3. Renewable Energy Expansion
• Solar and wind capacity additions
• Transmission infrastructure build-out
• Energy storage requirements

4. Manufacturing Renaissance
• PLI schemes attracting investment
• China+1 strategy benefiting India

5. Urbanization Trends
• 40% urbanization currently, growing to 50% by 2030
• Housing and construction demand

Government Policy Support:
National Mineral Policy 2019, mining reforms and auction regime, export restrictions protecting domestic supply.""",

        "demand_drivers": """Demand Drivers for Vedanta Ltd.

End-User Industries:

1. Construction & Infrastructure (35% of demand)
• Real estate development
• Road and highway construction
• Port and airport development

2. Automotive Sector (20% of demand)
• Passenger and commercial vehicles
• Two-wheelers and EV components

3. Electrical & Electronics (18% of demand)
• Power cables and wiring
• Consumer electronics

4. Packaging Industry (12% of demand)
• Beverage cans
• Food and pharmaceutical packaging

5. Other Industries (15% of demand)
• Aerospace and defense
• Industrial machinery

Growth Outlook by Segment:
• Construction: 8% current, 10% outlook
• Automotive: 12% current, 15% outlook
• Electronics: 15% current, 18% outlook""",

        "industry_risks": """Industry Risks

Regulatory & Policy Risks:
• Environmental Regulations: Stricter emission norms, water usage restrictions
• Government Policy Changes: Export duty variations, royalty rate changes
• Impact: HIGH | Likelihood: MEDIUM

Market Risks:
• Commodity Price Volatility: Global demand-supply dynamics, currency fluctuations
• Competition Intensity: New capacity additions, import competition
• Impact: HIGH | Likelihood: HIGH

Operational Risks:
• Resource Depletion: Mine life limitations, grade deterioration
• Labor and Social Issues: Union negotiations, community relations
• Impact: MEDIUM | Likelihood: MEDIUM

Mitigation Strategies:
• Diversified commodity portfolio
• Long-term contracts with customers
• Hedging strategies for currency and commodities
• Strong community engagement programs""",

        "summary_table": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/summary_table_example.png",
        "chart_custom": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/custom_chart_example.png",
        "chart_profit_loss": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/profit_loss_20260207_111704.png",
        "chart_balance_sheet": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/balance_sheet_20260207_111704.png",
        "chart_cash_flow": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/cash_flow_20260207_111705.png",
        "chart_ratio_analysis": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/ratios_20260207_111705.png",
        "chart_summary": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/custom_chart_20260207_111702.png",
    }

    # Get the directory of this script
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, "master_template.pptx")
    output_dir = os.path.join(script_dir, "output")

    print("=" * 60)
    print("PPT GENERATOR - Research Report Automation")
    print("=" * 60)
    print(f"\nTemplate: {template_path}")
    print(f"Output Directory: {output_dir}")

    try:
        output_file = generate_report_ppt(
            data=example_data,
            template_path=template_path,
            output_dir=output_dir
        )
        print(f"\n{'=' * 60}")
        print(f"SUCCESS! Report generated: {output_file}")
        print("=" * 60)
    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()