
from pptx import Presentation
import os
import re

def analyze_ppt(file_path):
    if not os.path.exists(file_path):
        print(f"Error: {file_path} not found.")
        return

    try:
        prs = Presentation(file_path)
        print(f"Analysis of: {file_path}")
        print(f"Total Slides: {len(prs.slides)}\n")

        for i, slide in enumerate(prs.slides):
            print(f"--- Slide {i+1} ---")
            placeholders_found = []
            shapes_found = []

            for shape in slide.shapes:
                shape_type = "Shape"
                if shape.has_text_frame:
                    shape_type = "TextFrame"
                    text = shape.text_frame.text
                    # Find {{...}}
                    matches = re.findall(r"\{\{(.*?)\}\}", text)
                    if matches:
                        placeholders_found.extend(matches)
                        print(f"  [Text] Found placeholders: {matches}")
                    
                    # Also looking for specific text content that might be headers
                    if len(text) < 50:
                         shapes_found.append(f"Text: '{text.strip()}'")

                elif shape.has_table:
                    shape_type = "Table"
                    # Check first cell
                    try:
                        first_cell = shape.table.cell(0,0).text_frame.text.strip()
                        matches = re.findall(r"\{\{(.*?)\}\}", first_cell)
                        if matches:
                            placeholders_found.extend(matches)
                            print(f"  [Table] Found placeholders in first cell: {matches}")
                        shapes_found.append(f"Table ({len(shape.table.rows)}x{len(shape.table.columns)})")
                    except:
                        pass
                
                elif shape.shape_type == 13: # Picture
                    shape_type = "Picture"
                    shapes_found.append("Picture")
                
                # print(f"  Found {shape_type}")

            if not placeholders_found:
                print("  (No placeholders found)")
            
            # Print specific shapes if interesting
            # for s in shapes_found:
            #     print(f"  - {s}")

    except Exception as e:
        print(f"Error checking PPT: {e}")

if __name__ == "__main__":
    analyze_ppt("master_template.pptx")
