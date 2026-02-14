from pptx import Presentation
import os

def list_placeholders():
    file_path = "master_template.pptx"
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return

    print(f"Opening {file_path}...")
    prs = Presentation(file_path)
    
    print("\n--- Placeholders Found ---")
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\nSlide {slide_idx + 1}:")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if "{{" in text:
                    print(f"  - Text: {text}")
            if shape.has_table:
                # Check first cell for placeholder or any cell
                first_cell = shape.table.cell(0,0).text_frame.text.strip()
                if "{{" in first_cell:
                     print(f"  - Table [0,0]: {first_cell}")

if __name__ == "__main__":
    list_placeholders()
